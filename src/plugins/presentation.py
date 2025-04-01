from typing import Annotated
from semantic_kernel.functions.kernel_function_decorator import kernel_function
from pptx import Presentation
import os
import requests
from io import BytesIO
from PIL import Image
import re

class PresentationPlugin:
    """The Presentation Plugin can be used to create presentations decks and slides."""

    @kernel_function(description="Create a presentation deck in a PDF format.")
    def create_presentation(self, 
                            title: Annotated[str, "The title of the presentation"],
                            subtitle: Annotated[str, "The subtitle of the presentation"],
                            content: Annotated[str, "The content of the decks"],
                            template: Annotated[str, "The presentation template from the list_templates method"] = 'default.pptx'
                            ) -> Annotated[str, "Create a presentation."]:

        # Check if there is a .pptx extension
        if not template.endswith('.pptx'):
            template += '.pptx'

        # If the template does not exist, use the default template
        if template not in self.list_templates():
            template = 'default.pptx'

        # Create a presentation object
        prs = Presentation(pptx='templates/' + template)

        # Remove title and subtitle from the slides content
        content = content.replace(title, '')
        content = content.replace(subtitle, '')

        # Split the content into slides
        slides_content = content.split("#")
        initial_slide = True
        
        for slide_content in slides_content:
            # The first slide must have a title and subtitle, no images
            if initial_slide:
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                slide.placeholders[1].text = subtitle
                initial_slide = False
                continue
            
            if len(slide_content) > 0:
                image_found = False
                title = slide_content.split('\n')[0]
                slide_content = slide_content.replace(title, '').strip()
                # Add the content for the other slides
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title

                # Check if the content has any URLs, if it does render them as images, not text
                #find images in the slide_content

                print(f"Original slide content: {slide_content}")

                lines = slide_content.split('\n')
                #create a regex that matches and excludes everything before "the dalle_prefix", and anything that comes after it until the ) is encountered, do not include the ) in the match:
                #EXAMPLE: ![abc123](https://dalleproduse.blob.core.windows.net/private/images/1b7e4f6d-7801-4bd8-b63f-b7df14da7e9c/generated_00.png?se=2025-04-01T18%3A28%3A24Z&sig=sc60RW5VcZKsbW2lnvIPXcrZUbzRBIYN4Io3fuLkveA%3D&ske=2025-04-05T21%3A58%3A55Z&skoid=09ba021e-c417-441c-b203-c81e5dcd7b7f&sks=b&skt=2025-03-29T21%3A58%3A55Z&sktid=33e01921-4d64-4f8c-a055-5bdaffd5e33d&skv=2020-10-02&sp=r&spr=https&sr=b&sv=2020-10-02)
                #match = https://dalleproduse.blob.core.windows.net/private/images/1b7e4f6d-7801-4bd8-b63f-b7df14da7e9c/generated_00.png?se=2025-04-01T18%3A28%3A24Z&sig=sc60RW5VcZKsbW2lnvIPXcrZUbzRBIYN4Io3fuLkveA%3D&ske=2025-04-05T21%3A58%3A55Z&skoid=09ba021e-c417-441c-b203-c81e5dcd7b7f&sks=b&skt=2025-03-29T21%3A58%3A55Z&sktid=33e01921-4d64-4f8c-a055-5bdaffd5e33d&skv=2020-10-02&sp=r&spr=https&sr=b&sv=2020-10-02
                #dalle_prefix = 'https:\/\/dalleproduse.blob.core.windows.net\/private\/images\/'
                #The old one - dalle_regex = re.compile(rf'\({dalle_prefix}.*?\)')        
                
                dalle_regex = r'!\[.*?\]\((.*?)\)'
        
                for line in lines:
                    #if the line has a match for the dalle_regex, then it is an image URL
                                   #find the first match in the slide_content
                    match = re.search(dalle_regex, slide_content)
                    if match:
                        image_url = match.group(1)
                        print(f"Found image URL: {image_url}")

                        #if there is a match, replace the entire line with a blank
                        slide_content = slide_content.replace(line, '')
                        
                        print(f"Adding image: {image_url} to slide")
                        # Add the image to the slide
                        response = requests.get(image_url)
                        if response.status_code != 200:
                            print(f"Failed to download image from {image_url}")
                            continue

                        img = BytesIO(response.content)
                        slide.shapes.add_picture(img, left=int(prs.slide_width/2), top=int(prs.slide_height/2)+20, width=int(prs.slide_width/2), height=int(prs.slide_height/2))
                    
                #add remaining text to the slide, the image URLs should be gone from the text now:
                slide.placeholders[1].text = slide_content
                if(image_found):
                    #if there was an image, move the text to the upper left of the image
                    slide.placeholders[1].left = int(prs.slide_width/2) + 50
                    slide.placeholders[1].top = int(prs.slide_height/2) + 50

        # Save the presentation
        output_path = "presentation.pptx"
        prs.save(output_path)

        return output_path

    @kernel_function(description="List available presentation templates (.ppt and .pptx).")
    def list_templates(self) -> Annotated[list[str], "Available .ppt and .pptx templates"]:
        templates_dir = "templates"
        templates = [
            file for file in os.listdir(templates_dir)
            if file.endswith(".ppt") or file.endswith(".pptx")
        ]
        return templates